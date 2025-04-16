import os
import json
import requests
import certifi
from dotenv import load_dotenv
import urllib3
from pptx import Presentation
from datetime import datetime

# Disable insecure request warnings
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Load environment variables from .env file
load_dotenv()

# Set SSL certificate environment variables before starting HTTPS connections
os.environ["SSL_CERT_FILE"] = certifi.where()
os.environ["REQUESTS_CA_BUNDLE"] = certifi.where()

# Retrieve configuration from environment variables
tenant_id = os.getenv("TENANT_ID")
client_id = os.getenv("CLIENT_ID")
client_secret = os.getenv("CLIENT_SECRET")
resource = os.getenv("RESOURCE")
deployment_name = os.getenv("DEPLOYMENT_NAME")
openai_api_base = os.getenv("OPENAI_API_BASE")
subscription_key = os.getenv("SUBSCRIPTION_KEY")

# Function to obtain an Azure AD access token
def get_access_token():
    token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    token_data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": resource + ".default"
    }
    token_headers = {
        "Content-Type": "application/x-www-form-urlencoded"
    }
    response = requests.post(token_url, data=token_data, headers=token_headers)
    response.raise_for_status()
    return response.json().get("access_token")

# Function to log errors to a file
def log_error_to_file(error_message, response_text=None):
    log_folder = "error_logs"
    if not os.path.exists(log_folder):
        os.makedirs(log_folder)
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    log_file = os.path.join(log_folder, f"error_{timestamp}.log")
    with open(log_file, "w", encoding="utf-8") as f:
        f.write(f"Error occurred at {timestamp}\n")
        f.write(f"Error message: {error_message}\n")
        if response_text:
            f.write(f"Response content:\n{response_text}\n")
    print(f"Error details logged to {log_file}")

# Class to handle OpenAI text generation requests
class OpenAITextGenerator:
    def __init__(self, api_base, deployment, access_token, subscription_key):
        self.api_base = api_base.rstrip("/")  # remove trailing slash if present
        self.deployment = deployment
        self.access_token = access_token
        self.subscription_key = subscription_key
        self.api_version = "2024-07-01-preview"
    
    def send_request(self, system_message, user_message):
        api_url = f"{self.api_base}/deployments/{self.deployment}/chat/completions?api-version={self.api_version}"
        headers = {
            "Authorization": f"Bearer {self.access_token}",
            "Ocp-Apim-Subscription-Key": self.subscription_key,
            "Content-Type": "application/json",
            "api-key": self.access_token
        }
        data = {
            "messages": [
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_message}
            ]
        }
        try:
            response = requests.post(api_url, headers=headers, json=data, verify=False)
            response.raise_for_status()  # Raise an HTTPError for bad responses (4xx and 5xx)
            result = response.json()
            return result['choices'][0]['message']['content']
        except requests.exceptions.HTTPError as http_err:
            print(f"HTTP error occurred: {http_err}")
            print(f"Response content: {response.text}")
            log_error_to_file(str(http_err), response.text)
        except requests.exceptions.RequestException as req_err:
            print(f"Request error occurred: {req_err}")
            log_error_to_file(str(req_err))
        except KeyError as key_err:
            print(f"Unexpected response format: {key_err}")
            log_error_to_file(str(key_err), response.text)
        except Exception as e:
            print(f"An unexpected error occurred: {e}")
            log_error_to_file(str(e))

# Function to list all ppt or pptx files in a given folder
def list_ppt_files(folder_path):
    """List all PowerPoint files (.ppt, .pptx, .pptm) in the given folder."""
    ppt_files = []
    for filename in os.listdir(folder_path):
        if filename.lower().endswith((".ppt", ".pptx", ".pptm")):
            ppt_files.append(filename)
    ppt_files.sort()
    return ppt_files

# Function to extract text, title, and slide number from a PowerPoint file
def extract_text_with_metadata_from_ppt(file_path):
    """Extract text, title, slide number, notes, file name, and slide link from a PowerPoint file."""
    slides_data = []
    presentation = Presentation(file_path)
    file_name = os.path.basename(file_path)  # Get the file name
    abs_path = os.path.abspath(file_path)    # Convert to absolute path

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = ""
        slide_title = None
        slide_notes = None

        # Extract text from slide shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
            if shape.has_text_frame and shape.text_frame.text and not slide_title:
                slide_title = shape.text_frame.text

        # Extract notes from the slide
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_notes = slide.notes_slide.notes_text_frame.text.strip()

        # Create a link to the slide (e.g., file:///<absolute_path>#slide=<slide_number>)
        slide_link = f"file:///{abs_path}#slide={slide_number}"

        slides_data.append({
            "file_name": file_name,
            "title": slide_title if slide_title else f"Slide {slide_number}",
            "slide_number": slide_number,
            "text": slide_text.strip(),
            "note": slide_notes if slide_notes else "",  # Add notes to the JSON
            "slide_link": slide_link
        })

    return slides_data

# Function to load management data
def load_management_data(file_path):
    """Load management data from a JSON file."""
    if os.path.exists(file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

# Function to save management data
def save_management_data(file_path, data):
    """Save management data to a JSON file."""
    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)

# Function to save conversation history
def save_conversation_history(conversation_folder, conversation_memory):
    """Save the conversation history to a timestamped JSON file."""
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    conversation_file = os.path.join(conversation_folder, f"conversation_history_{timestamp}.json")
    with open(conversation_file, "w", encoding="utf-8") as f:
        json.dump(conversation_memory, f, indent=4, ensure_ascii=False)
    print(f"Conversation history saved to {conversation_file}")

def load_pre_paper_prompt():
    """Load the content of pre_paper_prompt.txt from the system_prompt folder."""
    system_prompt_folder = os.path.join("C:\\", "python_scripts", "pptChat", "system_prompt")
    pre_paper_prompt_file = os.path.join(system_prompt_folder, "pre_paper_prompt.txt")
    if not os.path.exists(pre_paper_prompt_file):
        raise FileNotFoundError(f"Pre-paper prompt file not found: {pre_paper_prompt_file}")
    with open(pre_paper_prompt_file, "r", encoding="utf-8") as f:
        return f.read()

def main():
    ppt_folder = os.path.join("C:\\", "python_scripts", "pptChat", "ppt")
    ppt_json_folder = os.path.join("C:\\", "python_scripts", "pptChat", "ppt_json")
    management_folder = os.path.join("C:\\", "python_scripts", "pptChat", "text_extraction_management_files")
    conversation_folder = os.path.join("C:\\", "python_scripts", "pptChat", "conversation_history")

    if not os.path.exists(ppt_folder):
        print(f"Folder not found: {ppt_folder}")
        return
    if not os.path.exists(ppt_json_folder):
        os.makedirs(ppt_json_folder)
    if not os.path.exists(management_folder):
        os.makedirs(management_folder)
    if not os.path.exists(conversation_folder):
        os.makedirs(conversation_folder)

    management_file = os.path.join(management_folder, "ppt_management.json")
    if not os.path.exists(management_file):
        with open(management_file, "w", encoding="utf-8") as f:
            json.dump({}, f)  # Create an empty management file

    management_data = load_management_data(management_file)

    # Load the pre-paper prompt
    pre_paper_prompt = load_pre_paper_prompt()

    while True:
        print("\n")
        print("Select an option:")
        print("\n")
        print("1: New ppt file")
        print("2: Existing json files")
        print("3: Exit")
        print("\n")
        option = int(input("Enter your choice: "))

        if option == 1:
            # List PowerPoint files
            ppt_files = list_ppt_files(ppt_folder)
            if not ppt_files:
                print("No PowerPoint files found in the folder.")
                continue

            print("\n")
            print("Select a PowerPoint file by entering its number:")
            print("\n")
            for idx, filename in enumerate(ppt_files, start=1):
                print(f"{idx}: {filename}")
            print("\n")
            selected_num = int(input("Enter the number: "))
            selected_file = ppt_files[selected_num - 1]
            ppt_path = os.path.join(ppt_folder, selected_file)

            # Check if the file has been updated
            last_modified_time = os.path.getmtime(ppt_path)
            last_modified_str = datetime.fromtimestamp(last_modified_time).isoformat()

            if selected_file in management_data and management_data[selected_file] == last_modified_str:
                print("No new update found.")
                print("\n")
                json_filename = os.path.splitext(selected_file)[0] + ".json"
                output_file = os.path.join(ppt_json_folder, json_filename)
                with open(output_file, "r", encoding="utf-8") as f:
                    system_message = f.read()
            else:
                # Extract text with metadata from the selected PowerPoint file
                slides_data = extract_text_with_metadata_from_ppt(ppt_path)
                if not slides_data:
                    print("No text extracted from the selected PowerPoint file.")
                    print("\n")
                    continue

                # Convert the extracted data to JSON format
                slides_json = json.dumps(slides_data, indent=4, ensure_ascii=False)

                # Save the JSON to a file with the same name as the PowerPoint file
                json_filename = os.path.splitext(selected_file)[0] + ".json"
                output_file = os.path.join(ppt_json_folder, json_filename)
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(slides_json)
                print(f"Slides data saved to {output_file}")
                print("\n")

                # Update management data
                management_data[selected_file] = last_modified_str
                save_management_data(management_file, management_data)

                # Set the extracted JSON data as the system prompt
                system_message = slides_json

        elif option == 2:
            # List JSON files in the ppt_json folder
            json_files = [f for f in os.listdir(ppt_json_folder) if f.lower().endswith(".json")]
            if not json_files:
                print("No JSON files found in the folder.")
                print("\n")
                continue
            
            print("\n")
            print("Select a JSON file by entering its number:")
            print("\n")
            for idx, filename in enumerate(json_files, start=1):
                print(f"{idx}: {filename}")
            print(f"{len(json_files) + 1}: All JSON files")
            
            selected_num = int(input("Enter the number: "))

            if selected_num == len(json_files) + 1:
                # Combine all JSON files into one system message
                combined_data = []
                for json_file in json_files:
                    file_path = os.path.join(ppt_json_folder, json_file)
                    with open(file_path, "r", encoding="utf-8") as f:
                        combined_data.extend(json.load(f))
                system_message = json.dumps(combined_data, indent=4, ensure_ascii=False)
            else:
                selected_file = json_files[selected_num - 1]
                file_path = os.path.join(ppt_json_folder, selected_file)
                with open(file_path, "r", encoding="utf-8") as f:
                    system_message = f.read()

        elif option == 3:
            print("Exiting the program.")
            break

        else:
            print("Invalid option selected.")
            continue

        # Prepend the pre-paper prompt to the system message
        system_message = pre_paper_prompt + "\n\n" + system_message

        # Instantiate the text generator
        generator = OpenAITextGenerator(openai_api_base, deployment_name, get_access_token(), subscription_key)

        # Initialize memory to store the conversation
        conversation_memory = [
            {"role": "system", "content": system_message}
        ]

        # Chat with the extracted text as the system prompt
        while True:
            user_message = input("Enter your prompt (type 'exit' to quit): ")
            if user_message.strip().lower() == "exit":
                break

            # Add user message to memory
            conversation_memory.append({"role": "user", "content": user_message})

            # Send the entire conversation memory to the AI
            response = generator.send_request(
                system_message=system_message,
                user_message=json.dumps(conversation_memory, ensure_ascii=False)
            )

            # Add AI response to memory
            conversation_memory.append({"role": "assistant", "content": response})

            # Display the AI's response
            print("\nResponse:")
            print(response)
            print()

        # Save the conversation history
        save_conversation_history(conversation_folder, conversation_memory)

if __name__ == "__main__":
    main()