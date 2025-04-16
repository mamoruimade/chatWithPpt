import os
import json
import requests
import certifi
import urllib3
import time
from pptx import Presentation
from datetime import datetime
import streamlit as st
from dotenv import load_dotenv

# Disable insecure request warnings and load environment variables
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
load_dotenv()

# Set SSL certificate environment variables
os.environ["SSL_CERT_FILE"] = certifi.where()
os.environ["REQUESTS_CA_BUNDLE"] = certifi.where()

# Retrieve configuration from environment variables
tenant_id = os.getenv("TENANT_ID")
client_id = os.getenv("CLIENT_ID")
client_secret = os.getenv("CLIENT_SECRET")
resource = os.getenv("RESOURCE")
deployment_name = os.getenv("DEPLOYMENT_NAME")
openai_api_base = os.getenv("OPENAI_API_BASE")
subscription_key = os.getenv("SUBSCRIPTION_KEY")

# Function to obtain an Azure AD access token
def get_access_token():
    token_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
    token_data = {
        "grant_type": "client_credentials",
        "client_id": client_id,
        "client_secret": client_secret,
        "scope": resource + ".default"
    }
    token_headers = {"Content-Type": "application/x-www-form-urlencoded"}
    response = requests.post(token_url, data=token_data, headers=token_headers)
    response.raise_for_status()
    return response.json().get("access_token")

# Class to handle OpenAI text generation requests
class OpenAITextGenerator:
    def __init__(self, api_base, deployment, access_token, subscription_key):
        self.api_base = api_base.rstrip("/")  # Remove trailing slash if present
        self.deployment = deployment
        self.access_token = access_token
        self.subscription_key = subscription_key
        self.api_version = "2024-07-01-preview"
    
    def send_request(self, messages):
        api_url = f"{self.api_base}/deployments/{self.deployment}/chat/completions?api-version={self.api_version}"
        headers = {
            "Authorization": f"Bearer {self.access_token}",
            "Ocp-Apim-Subscription-Key": self.subscription_key,
            "Content-Type": "application/json",
        }
        data = {"messages": messages}
        try:
            response = requests.post(api_url, headers=headers, json=data, verify=False)
            response.raise_for_status()  # Raise an HTTPError for bad responses (4xx and 5xx)
            result = response.json()
            return result['choices'][0]['message']['content']
        except requests.exceptions.HTTPError as http_err:
            print(f"HTTP error occurred: {http_err}")
            print(f"Response content: {response.text}")
            raise
        except requests.exceptions.RequestException as req_err:
            print(f"Request error occurred: {req_err}")
            raise
        except KeyError as key_err:
            print(f"Unexpected response format: {key_err}")
            raise
        except Exception as e:
            print(f"An unexpected error occurred: {e}")
            raise

# Function to extract text and metadata from a PowerPoint file
def extract_text_with_metadata_from_ppt(file_path):
    slides_data = []
    presentation = Presentation(file_path)
    file_name = os.path.basename(file_path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = ""
        slide_title = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
            if shape.has_text_frame and shape.text_frame.text and not slide_title:
                slide_title = shape.text_frame.text

        slides_data.append({
            "file_name": file_name,
            "title": slide_title if slide_title else f"Slide {slide_number}",
            "slide_number": slide_number,
            "text": slide_text.strip()
        })

    return slides_data

# Force a rerun by setting a unique query parameter
def force_rerun():
    st.set_query_params(_=str(time.time()))

# Set page layout and title
st.set_page_config(layout="wide")
st.title("PPT Chat Application")

# Define folder paths and ensure they exist
ppt_folder = os.path.join("C:\\", "python_scripts", "pptChat", "ppt")
ppt_json_folder = os.path.join("C:\\", "python_scripts", "pptChat", "ppt_json")
conversation_folder = os.path.join("C:\\", "python_scripts", "pptChat", "conversation_history")
system_prompt_folder = os.path.join("C:\\", "python_scripts", "pptChat", "system_prompt")

os.makedirs(ppt_folder, exist_ok=True)
os.makedirs(ppt_json_folder, exist_ok=True)
os.makedirs(conversation_folder, exist_ok=True)
os.makedirs(system_prompt_folder, exist_ok=True)

# Sidebar: File upload and JSON file selection
with st.sidebar:
    st.header("File Management")
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["ppt", "pptx"])
    if uploaded_file:
        file_path = os.path.join(ppt_folder, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        st.success(f"Uploaded {uploaded_file.name}")

        with st.spinner("Extracting text from the PowerPoint file..."):
            slides_data = extract_text_with_metadata_from_ppt(file_path)
            json_filename = os.path.splitext(uploaded_file.name)[0] + ".json"
            json_path = os.path.join(ppt_json_folder, json_filename)
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(slides_data, f, indent=4, ensure_ascii=False)
        st.success(f"Extracted text saved to {json_filename}")

    json_files = [f for f in os.listdir(ppt_json_folder) if f.lower().endswith(".json")]
    selected_json = st.selectbox("Select a JSON file", json_files)

    # Add "Set" button to confirm the selected JSON file
    if st.button("Set"):
        if selected_json:
            # Read the pre_paper_prompt.txt content
            pre_paper_prompt_path = os.path.join(system_prompt_folder, "pre_paper_prompt.txt")
            if os.path.exists(pre_paper_prompt_path):
                with open(pre_paper_prompt_path, "r", encoding="utf-8") as f:
                    pre_paper_prompt = f.read()
            else:
                pre_paper_prompt = ""

            # Read the selected JSON file content
            with open(os.path.join(ppt_json_folder, selected_json), "r", encoding="utf-8") as f:
                json_content = f.read()

            # Combine pre_paper_prompt and JSON content
            system_message = pre_paper_prompt + "\n\n" + json_content

            # Update session state for conversation
            st.session_state.conversation = [{"role": "system", "content": system_message}]
            st.success(f"System message updated with content from {selected_json}")
        else:
            st.error("Please select a JSON file before clicking 'Set'.")

# Force a rerun by setting a unique query parameter
def force_rerun():
    st.experimental_set_query_params(_=str(time.time()))

# Main area: Chat Interface
# Force a rerun by setting a unique query parameter
def force_rerun():
    st.experimental_set_query_params(_=str(time.time()))

# Main area: Chat Interface
if "conversation" in st.session_state:
    # Display chat messages from history on app rerun
    for message in st.session_state.conversation:
        if message["role"] != "system":  # Skip displaying the system message
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    # Accept user input
    if user_prompt := st.chat_input("Type your message here..."):
        # Add user message to conversation history
        st.session_state.conversation.append({"role": "user", "content": user_prompt})
        # Display user message in chat message container
        with st.chat_message("user"):
            st.markdown(user_prompt)

        # Generate assistant response using OpenAI API
        generator = OpenAITextGenerator(
            openai_api_base,
            deployment_name,
            get_access_token(),
            subscription_key,
        )
        with st.chat_message("assistant"):
            with st.spinner("Waiting for response..."):
                try:
                    # Send the entire conversation history to the API
                    response = generator.send_request(st.session_state.conversation)
                    st.markdown(response)
                    # Add assistant response to conversation history
                    st.session_state.conversation.append({"role": "assistant", "content": response})
                except requests.exceptions.HTTPError as e:
                    st.error(f"Error: {e.response.text}")